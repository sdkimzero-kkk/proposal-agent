"""PPTX 파일 병합 유틸리티

여러 Phase별 PPTX 파일을 하나의 완성된 제안서 PPTX로 병합합니다.
슬라이드 XML + 관계(이미지, 차트) 복사를 모두 지원합니다.
"""

import copy
import os
from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

from .logger import get_logger

logger = get_logger("pptx_merger")

# 관계 속성 네임스페이스
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_REL_ATTRS = [
    f"{{{_R_NS}}}id",
    f"{{{_R_NS}}}embed",
    f"{{{_R_NS}}}link",
    f"{{{_R_NS}}}href",
]


# ──────────────────────────────────────────────────────────────
# Public API
# ──────────────────────────────────────────────────────────────

def merge_pptx_files(
    input_files: List[str | Path],
    output_path: str | Path,
    slide_width_inches: float = 13.333,
    slide_height_inches: float = 7.5,
) -> int:
    """여러 PPTX 파일을 하나로 병합

    각 Phase별 PPTX를 순서대로 병합하여 완성된 제안서를 생성합니다.
    텍스트, 도형, 배경, 이미지, 차트 모두 복사됩니다.

    Args:
        input_files: 병합할 PPTX 파일 경로 목록 (순서 중요)
        output_path: 출력 PPTX 파일 경로
        slide_width_inches: 슬라이드 가로 크기 (기본 13.333" = 16:9)
        slide_height_inches: 슬라이드 세로 크기 (기본 7.5" = 16:9)

    Returns:
        병합된 총 슬라이드 수
    """
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    result = Presentation()
    result.slide_width = Inches(slide_width_inches)
    result.slide_height = Inches(slide_height_inches)

    total_slides = 0
    processed_files = 0

    for file_path in input_files:
        file_path = Path(file_path)
        if not file_path.exists():
            logger.warning(f"파일 없음, 건너뜀: {file_path}")
            continue

        try:
            prs = Presentation(file_path)
            file_slide_count = len(prs.slides)

            for slide in prs.slides:
                _copy_slide(result, slide)
                total_slides += 1

            logger.info(f"  + {file_path.name}: {file_slide_count}장 (누적 {total_slides}장)")
            processed_files += 1

        except Exception as e:
            logger.error(f"파일 처리 실패 ({file_path.name}): {e}")

    result.save(str(output_path))
    logger.info(f"병합 완료: {processed_files}개 파일, 총 {total_slides}장 → {output_path}")
    return total_slides


def get_slide_count(pptx_path: str | Path) -> int:
    """PPTX 파일의 슬라이드 수 반환"""
    try:
        prs = Presentation(str(pptx_path))
        return len(prs.slides)
    except Exception as e:
        logger.error(f"슬라이드 수 조회 실패 ({pptx_path}): {e}")
        return 0


def calculate_page_offsets(
    phase_files: List[str | Path],
    cover_pages: int = 3,
) -> dict:
    """Phase별 시작 페이지 번호 계산

    각 Phase PPTX의 슬라이드 수를 세어 누적 시작 페이지를 반환합니다.
    표지/목차는 page number 미표시로 처리하므로 cover_pages만큼 오프셋.

    Args:
        phase_files: Phase별 PPTX 파일 경로 목록
        cover_pages: 표지/목차 슬라이드 수 (페이지 번호 제외)

    Returns:
        {파일경로: 시작_페이지_번호} 딕셔너리
        예: {"pptx/phase0.pptx": 3, "pptx/phase1.pptx": 11, ...}
    """
    offsets = {}
    current_page = cover_pages + 1  # 표지/목차 이후 첫 번호

    for file_path in phase_files:
        file_path = Path(file_path)
        offsets[str(file_path)] = current_page
        count = get_slide_count(file_path)
        if count > 0:
            current_page += count

    return offsets


# ──────────────────────────────────────────────────────────────
# Internal helpers
# ──────────────────────────────────────────────────────────────

def _copy_slide(dest_prs: Presentation, src_slide) -> None:
    """슬라이드 하나를 대상 프레젠테이션에 복사

    배경, 도형(텍스트/이미지/차트), 관계(rId) 모두 처리합니다.
    """
    # blank layout(index 6)으로 새 슬라이드 추가
    blank_layout = dest_prs.slide_layouts[6]
    new_slide = dest_prs.slides.add_slide(blank_layout)

    # 기본 placeholder 제거
    sp_tree = new_slide.shapes._spTree
    for el in list(sp_tree):
        sp_tree.remove(el)

    # ── 관계 복사 (이미지, 차트 등 embedded parts) ──
    rId_map = _copy_relationships(src_slide.part, new_slide.part)

    # ── shapes XML 복사 ──
    for el in src_slide.shapes._spTree:
        new_el = copy.deepcopy(el)
        _patch_rids(new_el, rId_map)
        sp_tree.append(new_el)

    # ── 배경 복사 ──
    _copy_background(src_slide, new_slide)


def _copy_relationships(src_part, dest_part) -> dict:
    """소스 슬라이드 파트의 관계를 대상 파트에 복사

    Returns:
        {구_rId: 새_rId} 매핑 딕셔너리
    """
    rId_map = {}
    for rId, rel in src_part.rels.items():
        if rel.is_external:
            # 외부 링크(하이퍼링크 등)
            try:
                new_rId = dest_part.target_ref  # fallback
                # pptx는 외부 관계를 add_relationship으로 추가
                new_rId = dest_part._rels._next_rId()
                dest_part._rels[new_rId] = _make_external_rel(
                    dest_part, new_rId, rel.reltype, rel.target_ref
                )
                rId_map[rId] = new_rId
            except Exception:
                pass
        else:
            # 내부 파트(이미지, 차트, 노트 등)
            try:
                new_rId = dest_part.relate_to(rel.target_part, rel.reltype)
                rId_map[rId] = new_rId
            except Exception:
                # 이미 관계가 있거나 복사 불가 — 기존 rId 검색
                for existing_rId, existing_rel in dest_part.rels.items():
                    if (not existing_rel.is_external and
                            existing_rel.target_part is rel.target_part):
                        rId_map[rId] = existing_rId
                        break
    return rId_map


def _patch_rids(element, rId_map: dict) -> None:
    """XML 요소 트리 내 관계 속성값을 새 rId로 교체"""
    if not rId_map:
        return
    for node in element.iter():
        for attr in _REL_ATTRS:
            val = node.get(attr)
            if val and val in rId_map:
                node.set(attr, rId_map[val])


def _copy_background(src_slide, dest_slide) -> None:
    """슬라이드 배경(p:bg) XML 복사"""
    src_cSld = src_slide._element.find(qn("p:cSld"))
    if src_cSld is None:
        return
    src_bg = src_cSld.find(qn("p:bg"))
    if src_bg is None:
        return

    dest_cSld = dest_slide._element.find(qn("p:cSld"))
    if dest_cSld is None:
        return

    # 기존 배경 제거 후 복사본 삽입
    for old_bg in dest_cSld.findall(qn("p:bg")):
        dest_cSld.remove(old_bg)
    dest_cSld.insert(0, copy.deepcopy(src_bg))


def _make_external_rel(part, rId, reltype, target_ref):
    """외부 관계 객체 생성 (pptx 내부 Relationship 모방)"""
    from pptx.opc.rel import Relationships, _Relationship
    return _Relationship(rId, reltype, target_ref, part._baseURI, is_external=True)


# ──────────────────────────────────────────────────────────────
# 편의 함수: merge_pptx.py 스크립트 템플릿 생성
# ──────────────────────────────────────────────────────────────

def generate_merge_script(
    output_folder: str | Path,
    phase_count: int = 8,
    output_filename: str = "제안서_완성.pptx",
) -> str:
    """merge_pptx.py 스크립트 내용 생성

    각 Phase-Split 생성 스크립트와 함께 저장할 병합 스크립트를 반환합니다.

    Args:
        output_folder: 생성 스크립트가 위치할 폴더 (output/테스트 XX/)
        phase_count: Phase 스크립트 수 (기본 8)
        output_filename: 최종 출력 파일명

    Returns:
        merge_pptx.py 파일 내용 문자열
    """
    phase_list = ", ".join(
        f'"pptx/phase{i}.pptx"' for i in range(phase_count)
    )

    return f'''#!/usr/bin/env python3
"""Phase별 PPTX 병합 스크립트 — 자동 생성됨"""
import sys, os

PROJECT_ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), "../.."))
sys.path.insert(0, PROJECT_ROOT)

from src.utils.pptx_merger import merge_pptx_files, get_slide_count

# ── 병합 순서: 표지/목차 → Phase 0 ~ {phase_count-1} ──
FILES = [
    "pptx/cover_toc.pptx",
    {phase_list},
]

OUTPUT = "{output_filename}"

def main():
    print("\\n[PPTX 병합 시작]")
    print(f"  입력 파일: {{len(FILES)}}개")

    for f in FILES:
        count = get_slide_count(f) if os.path.exists(f) else 0
        status = f"{{count}}장" if os.path.exists(f) else "없음"
        print(f"  - {{f}}: {{status}}")

    total = merge_pptx_files(FILES, OUTPUT)
    print(f"\\n병합 완료: 총 {{total}}장 → {{OUTPUT}}")


if __name__ == "__main__":
    main()
'''
